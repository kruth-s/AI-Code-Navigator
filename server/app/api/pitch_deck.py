from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List, Optional
import os
import json
import uuid
from pptx import Presentation
from pptx.util import Inches

from app.services.llm import get_llm
from langchain_core.messages import SystemMessage

router = APIRouter(prefix="/api/pitch", tags=["pitch-deck"])

class PitchRequest(BaseModel):
    repo_name: str

class PitchData(BaseModel):
    product_name: str
    tagline: str
    problem: str
    solution: str
    features: List[str]
    target_users: List[str]
    market_size: str
    business_model: str
    future_scope: str
    repo_link: str

class ExportRequest(BaseModel):
    data: PitchData

@router.post("/generate", response_model=PitchData)
async def generate_pitch(request: PitchRequest):
    repo_name = request.repo_name
    repos_dir = os.path.join(os.getcwd(), "repos", repo_name)
    
    code_samples = []
    total_files = 0
    if os.path.exists(repos_dir):
        for root, dirs, files in os.walk(repos_dir):
            if ".git" in root or "node_modules" in root or "dist" in root:
                continue
            for file in files:
                if file.endswith(('.js', '.ts', '.jsx', '.tsx', '.py', '.java', '.go', '.md', '.json')):
                    total_files += 1
                    try:
                        path = os.path.join(root, file)
                        with open(path, "r", encoding="utf-8") as f:
                            content = f.read()
                            if len(code_samples) < 10 and len(content) < 8000:
                                code_samples.append(f"File: {file}\n{content[:1500]}")
                    except:
                        pass

    vibe_context = "\n\n".join(code_samples)
    
    if total_files == 0:
        prompt = f"""
Analyze this GitHub repository name and convert it into a highly compelling SaaS/startup idea.
Repository Name: {repo_name}
(The local codebase is empty/uninitialized, so rely heavily on the repo name to imagine what this product could be).

Return ONLY valid JSON. Do not wrap it in markdown block quotes (no ```json). 
Use EXACTLY this schema structure:
{{
  "product_name": "CodeGuard AI",
  "tagline": "AI-powered code intelligence for startups",
  "problem": "- Lack of code visibility\\n- Late bug detection",
  "solution": "- Automated repo analysis\\n- Real-time insights",
  "features": ["PR risk analysis", "Bug prediction"],
  "target_users": ["Startups", "Developers"],
  "market_size": "Growing DevOps Market",
  "business_model": "SaaS subscription",
  "future_scope": "Enterprise integrations",
  "repo_link": "https://github.com/{repo_name}"
}}
"""
    else:
        prompt = f"""
Analyze this GitHub repository and convert it into a highly compelling SaaS/startup idea.
Repository Name: {repo_name}

Code Snippets:
{vibe_context}

Return ONLY valid JSON. Do not wrap it in markdown block quotes (no ```json). 
Use EXACTLY this schema structure:
{{
  "product_name": "CodeGuard AI",
  "tagline": "AI-powered code intelligence for startups",
  "problem": "- Lack of code visibility\\n- Late bug detection",
  "solution": "- Automated repo analysis\\n- Real-time insights",
  "features": ["PR risk analysis", "Bug prediction"],
  "target_users": ["Startups", "Developers"],
  "market_size": "Growing DevOps Market",
  "business_model": "SaaS subscription",
  "future_scope": "Enterprise integrations",
  "repo_link": "https://github.com/{repo_name}"
}}
"""

    try:
        llm = get_llm()
        response = llm.invoke([SystemMessage(content=prompt)])
        text = response.content.strip()
        
        # Clean markdown
        if text.startswith("```json"):
            text = text.replace("```json", "", 1)
        if text.endswith("```"):
            text = text[:-3]
        text = text.strip()
        
        data = json.loads(text)
        return PitchData(
            product_name=data.get("product_name", "AI Product"),
            tagline=data.get("tagline", "Turning code into revenue"),
            problem=data.get("problem", "- High costs\n- Low efficiency"),
            solution=data.get("solution", "- AI Automation\n- Seamless integration"),
            features=data.get("features", ["Smart Analytics"]),
            target_users=data.get("target_users", ["Developers"]),
            market_size=data.get("market_size", "TAM: $1B+"),
            business_model=data.get("business_model", "SaaS"),
            future_scope=data.get("future_scope", "Enterprise APIs"),
            repo_link=f"https://github.com/{repo_name}"
        )
    except Exception as e:
        print("Pitch Deck LLM Error:", e)
        raise HTTPException(status_code=500, detail="Failed to generate pitch data")


@router.post("/export")
async def export_pitch_deck(request: ExportRequest):
    data = request.data.dict()
    output_filename = f"pitch_deck_{uuid.uuid4().hex[:8]}.pptx"
    output_path = os.path.join(os.getcwd(), output_filename)
    
    try:
        prs = Presentation()

        def add_slide(title, content):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            try:
                # Add text to main body placeholder
                slide.placeholders[1].text = content
            except Exception:
                pass

        # Slide 1: Title
        add_slide(
            data["product_name"],
            f'{data["tagline"]}\n\nRepo: {data.get("repo_link", "")}'
        )

        # Slide 2: Problem + Solution
        add_slide(
            "Problem & Solution",
            f'Problem:\n{data["problem"]}\n\nSolution:\n{data["solution"]}'
        )

        # Slide 3: Features
        features = "\n".join([f"- {f}" for f in data["features"]])
        add_slide(
            "Product & Features",
            features
        )

        # Slide 4: Market + Users
        users = ", ".join(data["target_users"])
        add_slide(
            "Market & Users",
            f'Target Users: {users}\n\nMarket:\n{data["market_size"]}'
        )

        # Slide 5: Business + Future
        add_slide(
            "Business & Growth",
            f'Model: {data["business_model"]}\n\nFuture Scope:\n{data["future_scope"]}'
        )

        prs.save(output_path)
        
        return FileResponse(
            path=output_path, 
            filename="Auto_Pitch_Deck.pptx", 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            background=None # we don't need a background task for small file, or we handle cleanup separately
        )
    except Exception as e:
        print("Presentation creation error:", e)
        raise HTTPException(status_code=500, detail="Failed to generate pptx file")
